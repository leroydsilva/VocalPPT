from pptx import Presentation
import pyttsx3
import speech_recognition as sr
from pptx.util import Inches
from PIL import Image
import os
from pptx.util import Pt
from flask import session
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import datetime
from pptx.dml.color import RGBColor
from matplotlib.colors import to_rgb
# name='messi'

def color_to_rgb(color):
    """Converts a color to a RGB tuple from (0-255)."""
    if isinstance(color, tuple):
        # if a RGB tuple already
        return color
    else:
        # to_rgb() returns colors from (0-1)
        color = tuple(int(x * 255) for x in to_rgb(color))
        return color

def talk(audio):
    
    # time.sleep(2)
    # if audio==None:
    #     engine.say('good morning')
    try:
        engine=pyttsx3.init('sapi5')
        voices=engine.getProperty('voices')
        engine.setProperty('voice',voices[0])
        engine.setProperty('rate', 150)  # Speed percent (can go over 100)
        engine.setProperty('volume', 0.9)
        engine.say(audio)  
        engine.runAndWait()
    except:
        print('errpr in talnk function..................')   
    return 

def listen():
    

    # get audio from the microphone
    r = sr.Recognizer()
    with sr.Microphone() as source:
        # r.adjust_for_ambient_noise(source)
        # r.pause_threshold=1
        # r.energy_threshold=300
        print("Speak:")
        audio = r.listen(source)

    try:
        output = "" + r.recognize_google(audio)
    except sr.UnknownValueError:
        output = "Could not understand audio"
    except sr.RequestError as e:
        output = "Could not request results; {0}".format(e)

    data =output
    with open('transcript.txt', 'w') as f:
        f.write(data)
        f.close()
    return data.lower()

now = datetime.now()
x=f'{now.day}/{now.month}/{now.year}'

class CreatePpt(object):
    # global fileName
    global x
    # self.name=None
    def __init__(self,temp=None,slide_count=None,footer=None,name=None):
        self.pr1=Presentation(temp) 
        if slide_count==None:
            self.slide_count=0
        else:
            self.slide_count=slide_count   
        self.footer=footer
        self.name=name
        # self.pr1.save("static/{}.pptx".format(name))
        # os.system(f'ppt2pdf file static/{name}.pptx')

        #  name=session['phone']
        #  print(name+'thsi is naem asdorekokeo')
        #  self.pr1.save("static/{}.pptx".format(name)) 

    def make_slide(self,layout):
        self.slide=self.pr1.slides.add_slide(self.pr1.slide_layouts[layout])
        self.slide_count+=1
        self.shapes(data=self.slide_count,l=12.6,t=0.3)
        self.shapes(data=x,l=0.3,t=7)
        if self.footer!=None:
            self.shapes(data=self.footer,l=5.3,t=7)
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')
        print('done')
        self.l=[]
        self.ph=[]
        for s in self.slide.placeholders:
                self.l.append(s.name)
                self.ph.append(s.placeholder_format.idx)

        return self.l            

    def shapes(self,data,l,t):
        left =  Inches(l)
        top =Inches(t)
        width = height = Inches(3)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = str(data)

    def add_title(self,title):
        self.title1=self.slide.shapes.title
        self.title1.text=title
        if self.footer==None:
            self.footer=title
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')
        
    

    def add_subtitle(self,ph_num,subtitle=""):    
        self.subtitle1=self.slide.placeholders[self.ph[ph_num]]
        self.subtitle1.text=subtitle
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')
        print('done sucessfully')

    def add_text(self,ph_num,text_data=""):
        self.t=self.slide.placeholders[self.ph[ph_num]]
        self.text_frame = self.t.text_frame
        self.p=self.text_frame.add_paragraph()
        self.run = self.p.add_run()
        self.run.text=text_data
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')
        print("done text")

    def add_font(self,font_name=None,font_color=None,font_size=None,bold=False,italic=False):
        self.font = self.run.font
        if font_name!=None:
            self.font.name=font_name
        if font_size!=None:
            self.font.size=Pt(int(font_size))    
        if font_color!=None:    
            r,g,b=color_to_rgb(font_color)
            self.font.color.rgb = RGBColor(r,g,b)
        self.font.bold = bold
        self.font.italic= italic    
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')

    def add_chart(self,data,num_data,ph_num):
        self.t=self.slide.placeholders[self.ph[ph_num]]
        self.chart_data = CategoryChartData()
        self.chart_data.categories = data
        self.num_arr=[int(x) for x in num_data]
        self.chart_data.add_series('Series 1', self.num_arr)
        self.t = self.t.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,self.chart_data)
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')



    def add_image(self, ph_num, image_url):
        self.placeholder = self.slide.placeholders[self.ph[ph_num]]

        # Calculate the image size of the image
        im = Image.open(image_url)
        width, height = im.size

        # Make sure the placeholder doesn't zoom in
        self.placeholder.height = height
        self.placeholder.width = width

        # Insert the picture
        self.placeholder = self.placeholder.insert_picture(image_url)

        # Calculate ratios and compare
        image_ratio = width / height
        self.placeholder_ratio = self.placeholder.width / self.placeholder.height
        ratio_difference = self.placeholder_ratio - image_ratio

        # Placeholder width too wide:
        if ratio_difference > 0:
            difference_on_each_side = ratio_difference / 2
            self.placeholder.crop_left = -difference_on_each_side
            self.placeholder.crop_right = -difference_on_each_side
        # Placeholder height too high
        else:
            difference_on_each_side = -ratio_difference / 2
            self.placeholder.crop_bottom = -difference_on_each_side
            self.placeholder.crop_top = -difference_on_each_side
        self.pr1.save("static/{}.pptx".format(self.name))
        os.system(f'ppt2pdf file static/{self.name}.pptx')

    
       
